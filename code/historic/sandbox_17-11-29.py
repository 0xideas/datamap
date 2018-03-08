
# coding: utf-8

# ### Code Marouan

# In[133]:


import os
import platform
import psutil
import pandas as pd


def kern():
    """Platform independent retrieval of systeminfo
    
    Note do we need the other info, because I don't have those parameters set on my device?"""
    return platform.uname()

def mpt():
    """returns a dict of all disks/drives/devices connected"""
    lbl_drives = ['device','mountpoint','fstype']
    disks = [d[0:3] for d in psutil.disk_partitions()]
    drives = [dict(zip(lbl_drives,ds)) for ds in disks]
    return [d['mountpoint']for d in drives]

def find_files(basepath):
    """Returns all files found on a mountpoint"""
    allfiles = [os.path.join(r, name) for r,d,f, in os.walk(basepath) for name in f]
    exts = [os.path.splitext(file)[1][1:] for file in allfiles]
    df = pd.DataFrame([allfiles, exts]).T
    df.columns = ['fpath','ext']
    df['ext'] = df['ext'].astype('category')
    # df.ext.value_counts().nlargest(10)
    return df

# base = r'C:\Users\Thomas\Google Drive\Projects\Data privacy\projects\data mapping'
# find_files(base)


# In[179]:


import os
from unidecode import unidecode

def find_files(path, extension='.txt'):
    filepaths = []
    for root, dirs, files in os.walk(path):
        for file in files:
            if '~' not in file:
                filepaths.append(root+ '/' + file)
    return filepaths

def get_ext(file):
    return os.path.splitext(file)[1][1:]


def extract_txt(fpath):
    with open(fpath, 'r', encoding='utf-8') as f:
        text = f.read()
    return text

def read_files(filelist):
    """
    Note: unidecode replaces non-ASCII with ASCII characters
    """
    text = ''
    ignored = []
    for file in filelist:
        print(file)
        try:
            text = text +'\t'+ open_txt(file)
            print(text)
        except:
            ignored.append(file)
            print('ignored: ',file)
    return unidecode(text), ignored

def write_file(fpath, text):
    with open(fpath, "w") as f:
        f.write(text)


# In[180]:


"""Functions to scan text"""

import re

#takes a list of strings and returns a list of strings that contains only the numbers in the input strings
reglib = {
        'number': '[0-9]*',
        'ip': '^(?:(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)\.){3}(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)$',
        'email': r'\b(?:[a-zA-Z0-9_\-\.]+)@(?:[a-zA-Z0-9_\-\.]+)\.(?:[a-zA-Z]{2,5})\b',
        'phone': r'(?:\b0|\+)0?\d{1,3}[ ]?\d{1,3}[ ]?\d{2,3}(?:[ ]?\d{2}){1,2}\b',
        'passport': r'\b[A-Z]{2}\d{6}\b',
        'iid': r'\b(\d{3}-?\d{6}<?\d{1}-?\d{2,3})\b',
        'IBAN': r'\b(?:[A-Za-z]{2}[0-9]{2})?(?:[ ]?[0-9]{4}){3}\b',
        'bank account number': r'\b[0-9](?:[ ]?[0-9]{4})(?:[ ]?[0-9]{2})\b',
        'creditcards': r'\b(?:\d{4}[\s_-]?){4}\b',
        'gender': r'\bMale|male|Female|female\b'
}

def get_matches(pattern, string):
    """Return list of matches based on regex pattern"""
    return [s for s in re.findall(pattern, string) if s!='']

def screen_text(pattern, text):
    """Return list of matches based on regex pattern"""
    return get_matches(pattern, text)

def get_pers_data(text):
    """Create a dictionary with a list of matches for each pattern"""
    match_dict = {}
    for p in reglib:
        matches = screen_text(reglib[p], text)
        match_dict[p] = [m for m in matches if m]
    return match_dict


# In[181]:


base = r'C:\Users\Thomas\Google Drive\Projects\Data privacy\projects\data mapping\repo\sample'
files = find_files(base)
for f in files:
    print(f, get_ext(f))


# #### PPTX

# In[182]:


from pptx import Presentation

def open_pptx(fpath):
    return Presentation(fpath)

def pptx2txt(prs):
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            text = text + shape.text.replace('\n',' ')
    return text

def extract_pptx(fpath):
    return pptx2txt(open_pptx(fpath))


# #### Excel

# In[199]:


def open_xlsx(fpath):
    """Read xlsx path as pandas Excel file"""
    return pd.ExcelFile(fpath)
    
def xlsx2txt(xl):
    """Access excel sheet and extract text for each column"""
    text = ''
    for sheet in xl.sheet_names:
        df = xl.parse(sheet).astype(str)
        for col in df.columns:
            coltext = ','.join(df[col].tolist())
            text = text+','+coltext
    return text

def extract_xlsx(fpath):
    return xlsx2txt(open_xlsx(fpath))

file = r'C:\Users\Thomas\Google Drive\Projects\Data privacy\projects\data mapping\repo\sample/persondb.xlsx'
extract_xlsx(file)


# #### PDF

# In[200]:


from PyPDF2 import PdfFileWriter, PdfFileReader
import re


# In[201]:


def open_pdf(filename):
    """Read pdf and return pdf reader object"""
    pdfFileObj = open(filename,'rb')
    return PdfFileReader(pdfFileObj)

def pdf2txt(pdfFileObj):
    """Convert pdf file reader object to text

    This function starts by initializing a string to which text is added
    page by page. Empty spaces or pdf metadata are cleaned at the end
    """
    text = ''
    nrpages = pdfFileObj.getNumPages()
    for p in range(nrpages):
        page = pdfFileObj.getPage(p).extractText()
        text += page # will add text cleaning steps later
    return text

def extract_pdf(fpath):
    return pdf2txt(read_pdf(fpath))

file = r'C:\Users\Thomas\Google Drive\Projects\Data privacy\projects\data mapping\repo\sample/samplepdf.pdf'
extract_pdf(file)


# ### Putting it all together

# In[238]:


fdict = {
        'txt': extract_txt,
        'pptx': extract_pptx,
        'pdf': extract_pdf,
        'xlsx': extract_xlsx
    }
    
def open_file(fpath, fdict=fdict):
    """Convert the string instruction to a function call"""
    return fdict[get_ext(fpath)](fpath)


def screen_files(filelist):
    df = pd.DataFrame()
    for f in files:
        print('screening:', f, 'of type', get_ext(f))
        if get_ext(f) in fdict:
            df[f] = pd.Series(get_pers_data(open_file(f)))
    return df.T

df = screen_files(find_files(base))

def write_xlsx(df, fpath, sheetname='Sheet1'):
    writer = pd.ExcelWriter(fpath)
    df.to_excel(writer, sheetname)
    writer.save()
    return print('Excel file saved to', fpath)

# write_xlsx(df, 'output.xlsx')


def count_matches(df):
    """Quickly count the number of records per field"""
    return df.applymap(lambda x: len(x))

count_matches(df)

